import os
import re
import urllib.request
import urllib.parse
import tempfile
import asyncio
import httpx
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
from app.services.openai_service import generate_json_response

# ─── Design tokens ───
BG_DARK       = RGBColor(0x0B, 0x0D, 0x17)   # deep navy
BG_CARD       = RGBColor(0x12, 0x14, 0x22)   # card surface
BG_GRAD_END   = RGBColor(0x1A, 0x10, 0x3A)   # dark purple (gradient end)
ACCENT        = RGBColor(0x63, 0x66, 0xF1)   # indigo-500
ACCENT_LIGHT  = RGBColor(0x81, 0x84, 0xF5)   # lighter indigo
ACCENT_GLOW   = RGBColor(0x78, 0x7B, 0xF7)   # soft glow indigo
TEXT_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY      = RGBColor(0xC0, 0xC5, 0xD0)   # brighter gray for readability
TEXT_MUTED     = RGBColor(0x6B, 0x72, 0x80)   # gray-500
PURPLE         = RGBColor(0xA8, 0x55, 0xF7)   # violet-500
EMERALD        = RGBColor(0x10, 0xB9, 0x81)   # emerald-500

SLIDE_WIDTH  = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

FONT_HEADING  = "Segoe UI"
FONT_BODY     = "Segoe UI"


def _hex(color):
    """Convert RGBColor to hex string (no #)."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _set_slide_bg(slide, color):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_gradient_bg(slide, color1=BG_DARK, color2=BG_GRAD_END, angle=270):
    """Set a linear gradient background on a slide via direct XML."""
    ang = angle * 60000  # OOXML uses 60,000ths of a degree
    c1, c2 = _hex(color1), _hex(color2)
    cSld = slide.background._element if hasattr(slide.background, '_element') else slide.background._cSld
    
    # 1. Get or create p:bg
    bg = cSld.find(qn('p:bg'))
    if bg is None:
        bg = parse_xml(f'<p:bg {nsdecls("p")}/>')
        cSld.insert(0, bg)
        
    # 2. Get or create p:bgPr
    bgPr = bg.find(qn('p:bgPr'))
    if bgPr is None:
        bgPr = parse_xml(f'<p:bgPr {nsdecls("p")}/>')
        bg.append(bgPr)
        
    # 3. Remove existing fill types
    fill_tags = [
        qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'), 
        qn('a:pattFill'), qn('a:grpFill'), qn('a:blipFill')
    ]
    for tag in fill_tags:
        el = bgPr.find(tag)
        if el is not None:
            bgPr.remove(el)
            
    # 4. Create and append gradFill
    fill_xml = (
        f'<a:gradFill {nsdecls("a")} rotWithShape="0">'
        f'  <a:gsLst>'
        f'    <a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
        f'    <a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
        f'  </a:gsLst>'
        f'  <a:lin ang="{ang}" scaled="0"/>'
        f'</a:gradFill>'
    )
    grad_elem = parse_xml(fill_xml)
    bgPr.append(grad_elem)


def _add_shape_rect(slide, left, top, width, height, fill_color, opacity=1.0):
    """Add a filled rectangle (used for accent bars, cards, etc.)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if opacity < 1.0:
        _set_shape_transparency(shape, int((1.0 - opacity) * 100000))
    return shape


def _add_rounded_card(slide, left, top, width, height, fill_color, corner_radius=0.08):
    """Add a rounded-rectangle card shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if len(shape.adjustments) > 0:
        shape.adjustments[0] = corner_radius
    return shape


def _set_shape_transparency(shape, alpha_val):
    """Set transparency on a shape. alpha_val: 0=opaque, 100000=fully transparent."""
    solid_fill = shape.fill._fill
    srgb = solid_fill.find(qn('a:solidFill'))
    if srgb is None:
        return
    clr = srgb.find(qn('a:srgbClr'))
    if clr is None:
        return
    alpha_elem = parse_xml(f'<a:alpha {nsdecls("a")} val="{alpha_val}"/>')
    clr.append(alpha_elem)


def _add_accent_circle(slide, left, top, diameter, color):
    """Add a small decorative filled circle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _set_text(text_frame, text, font_size, color, bold=False, font_name=FONT_BODY, align=PP_ALIGN.LEFT):
    """Set formatted text on the first paragraph of a text frame."""
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align


def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, font_name=FONT_BODY, align=PP_ALIGN.LEFT):
    """Add a text box with styled text and zero margins."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    _set_text(tf, text, font_size, color, bold, font_name, align)
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_GRAY):
    """Add a text box containing multiple bullet points with proper hanging indents."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"▸\t{item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(12)  # breathing room between points
        p.alignment = PP_ALIGN.LEFT
        
        # Proper hanging indent:
        p.left_indent = Inches(0.35)
        p.first_line_indent = Inches(-0.25)

    return txBox


def _build_title_slide(prs, title_text, subtitle_text=""):
    """Slide 0: Gradient title slide with decorative accent elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END, angle=315)

    # Left accent bar
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT)

    # Decorative accent circle (top-right)
    _add_accent_circle(slide, Inches(10.5), Inches(0.6), Inches(1.2), ACCENT_GLOW)
    # Smaller companion circle
    _add_accent_circle(slide, Inches(11.9), Inches(1.4), Inches(0.5), PURPLE)

    # Title
    _add_text_box(
        slide, Inches(1.2), Inches(2.0), Inches(10.5), Inches(1.8),
        title_text, font_size=42, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING
    )

    # Subtitle
    if subtitle_text:
        _add_text_box(
            slide, Inches(1.2), Inches(4.2), Inches(10.5), Inches(0.8),
            subtitle_text, font_size=16, color=TEXT_MUTED
        )

    # Bottom accent line
    _add_shape_rect(slide, Inches(1.2), Inches(6.7), Inches(3), Inches(0.04), ACCENT)

    # Branding
    _add_text_box(
        slide, Inches(1.2), Inches(6.85), Inches(4), Inches(0.5),
        "Generated by LockIn", font_size=11, color=TEXT_MUTED
    )


def _build_content_slide(prs, slide_data, slide_num, total_slides):
    """Standard content slide with gradient bg and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END)

    title = slide_data.get("title", "Untitled")
    points = slide_data.get("points", [])

    # Top accent bar
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

    # Slide number badge
    _add_text_box(
        slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4),
        f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_MUTED
    )

    # Title
    _add_text_box(
        slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
        title, font_size=28, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING
    )

    # Accent divider
    _add_shape_rect(slide, Inches(0.8), Inches(1.75), Inches(2.0), Inches(0.04), ACCENT_LIGHT)

    # Bullet points
    if points:
        _add_bullet_list(
            slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.5),
            points, font_size=15, color=TEXT_GRAY
        )

    # Bottom branding bar
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(
        slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35),
        "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED
    )


async def _fetch_image_for_slide_async(client: httpx.AsyncClient, image_prompt: str, max_retries: int = 2) -> str:
    """Fetch an image from pollinations.ai with retry logic for rate limits."""
    encoded_prompt = urllib.parse.quote(image_prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=800&nologo=true"
    
    for attempt in range(max_retries):
        try:
            response = await client.get(url, timeout=12.0, follow_redirects=True)
            if response.status_code == 200 and len(response.content) > 1000:
                fd, path = tempfile.mkstemp(suffix=".jpg")
                with os.fdopen(fd, 'wb') as f:
                    f.write(response.content)
                return path
            elif response.status_code == 429:
                wait = 2 * (attempt + 1)
                print(f"[PPT Image] Rate limited (429), waiting {wait}s before retry {attempt+1}/{max_retries}")
                await asyncio.sleep(wait)
                continue
            elif response.status_code >= 500:
                wait = 2 * (attempt + 1)
                print(f"[PPT Image] Server error ({response.status_code}), waiting {wait}s before retry {attempt+1}/{max_retries}")
                await asyncio.sleep(wait)
                continue
            else:
                print(f"[PPT Image] Unexpected status {response.status_code}")
                return None
        except Exception as e:
            print(f"[PPT Image] Attempt {attempt+1} failed: {e}")
            if attempt < max_retries - 1:
                await asyncio.sleep(2 * (attempt + 1))
            continue
    
    print(f"[PPT Image] All {max_retries} attempts exhausted for prompt: {image_prompt[:60]}...")
    return None


def _build_visual_slide(prs, slide_data, slide_num, total_slides, image_path=None):
    """Slide with content on the left in a card and a generated image on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END)

    title = slide_data.get("title", "Untitled")
    points = slide_data.get("points", [])

    # Top accent bar & Slide number
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    _add_text_box(slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4), f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_MUTED)

    # Title
    _add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8), title, font_size=28, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING)
    _add_shape_rect(slide, Inches(0.8), Inches(1.75), Inches(2.0), Inches(0.04), ACCENT_LIGHT)

    # Text Card (Left side)
    text_width = Inches(5.6)
    if not image_path:
        text_width = Inches(11.7) # Span fully if image fails
    
    _add_rounded_card(slide, Inches(0.8), Inches(2.0), text_width, Inches(4.7), BG_CARD, corner_radius=0.04)
    if points:
        _add_bullet_list(slide, Inches(1.1), Inches(2.3), text_width - Inches(0.6), Inches(4.1), points, font_size=14, color=TEXT_GRAY)

    # Image Card (Right side)
    if image_path:
        try:
            # Framing rounded card behind the image
            _add_rounded_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.7), BG_CARD, corner_radius=0.04)
            # Add image slightly smaller for a clean bordered padding look
            slide.shapes.add_picture(image_path, Inches(6.95), Inches(2.15), width=Inches(5.4), height=Inches(4.4))
        except Exception as e:
            print(f"Error adding picture to slide: {e}")
        finally:
            if os.path.exists(image_path):
                try:
                    os.remove(image_path)
                except Exception:
                    pass

    # Bottom branding
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35), "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED)


def _build_comparison_slide(prs, slide_data, slide_num, total_slides):
    """Two-column comparison slide using rounded cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END)

    title = slide_data.get("title", "Comparison")
    left_points = slide_data.get("left_points", [])
    right_points = slide_data.get("right_points", [])

    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    _add_text_box(slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4), f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_MUTED)

    # Centered Title
    _add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8), title, font_size=28, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING, align=PP_ALIGN.CENTER)
    _add_shape_rect(slide, Inches(5.66), Inches(1.75), Inches(2.0), Inches(0.04), ACCENT_LIGHT)

    # Left Column Card
    _add_rounded_card(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.7), BG_CARD, corner_radius=0.04)
    # Visual top-accent line on left card
    _add_shape_rect(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.1), ACCENT)
    if left_points:
        _add_bullet_list(slide, Inches(1.1), Inches(2.4), Inches(5.0), Inches(4.1), left_points, font_size=14, color=TEXT_GRAY)

    # Right Column Card
    _add_rounded_card(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.7), BG_CARD, corner_radius=0.04)
    # Visual top-accent line on right card
    _add_shape_rect(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.1), PURPLE)
    if right_points:
        _add_bullet_list(slide, Inches(7.2), Inches(2.4), Inches(5.0), Inches(4.1), right_points, font_size=14, color=TEXT_GRAY)

    # Bottom branding
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35), "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED)


def _build_quote_slide(prs, slide_data, slide_num, total_slides):
    """Full slide large quote with purple gradient and decorative shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, RGBColor(0x3B, 0x0F, 0x48), angle=135) # Deep purple gradient

    quote = slide_data.get("quote", "")
    author = slide_data.get("author", "")

    # Top accent bar & Slide number
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PURPLE)
    _add_text_box(slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4), f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_MUTED)

    # Decorative elements
    _add_accent_circle(slide, Inches(1.0), Inches(1.5), Inches(1.5), ACCENT_GLOW)
    
    # Rounded Card overlay with slight transparency (simulated via color blending)
    _add_rounded_card(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0), BG_CARD, corner_radius=0.04)

    # Big Quote mark
    _add_text_box(slide, Inches(1.6), Inches(1.7), Inches(2), Inches(1.5), "“", font_size=110, color=PURPLE, bold=True, font_name="Georgia")

    # The quote
    _add_text_box(slide, Inches(2.2), Inches(2.7), Inches(8.9), Inches(2.8), f"{quote}", font_size=28, color=TEXT_WHITE, bold=False, font_name=FONT_HEADING, align=PP_ALIGN.LEFT)

    # Author
    _add_text_box(slide, Inches(2.2), Inches(5.6), Inches(8.9), Inches(0.8), f"— {author}", font_size=18, color=ACCENT_LIGHT, bold=True, font_name=FONT_BODY, align=PP_ALIGN.LEFT)

    # Bottom branding
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35), "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED)


def _build_hero_slide(prs, slide_data, slide_num, total_slides, image_path=None):
    """Full-bleed background image layout with a semi-transparent dark overlay for high visual impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide_data.get("title", "Hero Slide")
    subtitle = slide_data.get("subtitle", "")

    if image_path:
        try:
            # 1. Full-bleed background image
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=SLIDE_WIDTH, height=SLIDE_HEIGHT)
            # 2. Semi-transparent dark overlay for text readability
            overlay = _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, BG_DARK, opacity=0.7)
        except Exception as e:
            print(f"Error adding background picture to hero slide: {e}")
            _set_gradient_bg(slide, BG_DARK, BG_GRAD_END, angle=225)
        finally:
            if os.path.exists(image_path):
                try:
                    os.remove(image_path)
                except Exception:
                    pass
    else:
        # Fallback to premium gradient if image generation failed
        _set_gradient_bg(slide, BG_DARK, BG_GRAD_END, angle=225)
        _add_accent_circle(slide, Inches(8.5), Inches(2.0), Inches(3.0), ACCENT_GLOW)

    # Top accent bar & Slide number
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    _add_text_box(slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4), f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_WHITE)

    # Left decorative card containing text
    _add_rounded_card(slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.5), BG_CARD, corner_radius=0.04)
    # Add alpha overlay effect inside card
    card_overlay = _add_shape_rect(slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.5), BG_DARK, opacity=0.3)

    # Content
    _add_text_box(slide, Inches(1.8), Inches(2.6), Inches(9.7), Inches(1.5), title, font_size=36, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING)
    _add_shape_rect(slide, Inches(1.8), Inches(4.1), Inches(3.0), Inches(0.04), ACCENT_LIGHT)

    if subtitle:
        _add_text_box(slide, Inches(1.8), Inches(4.4), Inches(9.7), Inches(1.5), subtitle, font_size=16, color=TEXT_GRAY)

    # Bottom branding
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35), "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED)


def _build_stats_slide(prs, slide_data, slide_num, total_slides):
    """Slide with title and horizontal cards highlighting key statistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END)

    title = slide_data.get("title", "Key Statistics")
    stats = slide_data.get("stats", [])

    # Top accent bar & Slide number
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    _add_text_box(slide, Inches(0.8), Inches(0.45), Inches(1.2), Inches(0.4), f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=TEXT_MUTED)

    # Title
    _add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8), title, font_size=28, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING)
    _add_shape_rect(slide, Inches(0.8), Inches(1.75), Inches(2.0), Inches(0.04), ACCENT_LIGHT)

    # Render horizontal stats cards
    if stats:
        num_stats = min(len(stats), 4)
        card_width = Inches(2.6)
        card_height = Inches(3.8)
        gap = Inches(0.4)
        total_width = (num_stats * card_width) + ((num_stats - 1) * gap)
        start_left = (SLIDE_WIDTH - total_width) / 2

        for i in range(num_stats):
            stat_item = stats[i]
            val = stat_item.get("value", "0")
            lbl = stat_item.get("label", "")
            
            left_pos = start_left + (i * (card_width + gap))
            # Card background
            _add_rounded_card(slide, left_pos, Inches(2.2), card_width, card_height, BG_CARD, corner_radius=0.08)
            # Soft accent border/top line
            _add_shape_rect(slide, left_pos, Inches(2.2), card_width, Inches(0.08), PURPLE if i % 2 == 0 else ACCENT)

            # Metric Value
            _add_text_box(slide, left_pos + Inches(0.1), Inches(3.0), card_width - Inches(0.2), Inches(1.2), val, font_size=40, color=TEXT_WHITE, bold=True, font_name=FONT_HEADING, align=PP_ALIGN.CENTER)
            # Label
            _add_text_box(slide, left_pos + Inches(0.2), Inches(4.3), card_width - Inches(0.4), Inches(1.3), lbl, font_size=14, color=TEXT_GRAY, bold=False, font_name=FONT_BODY, align=PP_ALIGN.CENTER)

    # Bottom branding
    _add_rounded_card(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35), BG_CARD)
    _add_text_box(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.35), "LockIn AI — Student Productivity Platform", font_size=9, color=TEXT_MUTED)


def _build_closing_slide(prs, title_text):
    """Final thank-you / closing slide with decorative elements and gradient background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, BG_DARK, BG_GRAD_END, angle=135)

    _add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, PURPLE)
    
    # Decorative accent circles
    _add_accent_circle(slide, Inches(10.0), Inches(4.5), Inches(2.0), ACCENT_GLOW)
    _add_accent_circle(slide, Inches(8.5), Inches(5.5), Inches(0.8), PURPLE)

    _add_text_box(
        slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
        "Thank You", font_size=44, color=TEXT_WHITE, bold=True,
        font_name=FONT_HEADING, align=PP_ALIGN.LEFT
    )

    _add_text_box(
        slide, Inches(1.5), Inches(3.9), Inches(8), Inches(0.6),
        f"Presentation on: {title_text}", font_size=16, color=TEXT_MUTED
    )

    _add_shape_rect(slide, Inches(1.5), Inches(6.7), Inches(3), Inches(0.04), PURPLE)
    _add_text_box(
        slide, Inches(1.5), Inches(6.85), Inches(5), Inches(0.4),
        "Powered by LockIn AI", font_size=11, color=TEXT_MUTED
    )


async def generate_ppt_file(prompt_text: str, num_slides: int, user_id: str = "anonymous", model="openai/gpt-4o-mini") -> str:
    """
    Calls OpenAI to get structured PPT contents, compiles them into a
    professionally styled PowerPoint file and returns the local filepath.
    """
    system_prompt = f'''You are an expert presentation designer creating a professional, university-level presentation.
Topic: "{prompt_text}"
Total slides: EXACTLY {num_slides}

RULES:
1. Use a MIX of these slide types (at least 2 different types):
   - "content": Title + 4-6 bullet points. Each bullet is 1-2 clear sentences with real substance.
   - "visual": Title + 3-4 bullet points + an "image_prompt". The image_prompt must be a SHORT, simple description (under 15 words, e.g. "modern office workspace with dual monitors and coffee").
   - "comparison": Title + "left_points" + "right_points" (3-4 items each). Good for pros/cons, before/after, or two approaches.
   - "quote": A memorable "quote" + "author". Use sparingly (max 1).
   - "hero": Title + Subtitle + an "image_prompt". Perfect for key transitions or section dividers. The image is displayed full-bleed in the background. Use at most 1 hero slide.
   - "stats": Title + a "stats" list (3-4 items). Each stat has a "value" (e.g., "98%", "10x", "50M") and a "label" explaining it. Good for showing metrics, performance, growth, or numbers.
2. Bullet points must be substantive explanations, not vague phrases.
3. Titles must be concise (3-7 words).
4. image_prompt must be SHORT and concrete — no flowery language. Example: "student studying with laptop in library" NOT "A cinematic, breathtakingly detailed ultra-wide photograph of..."
5. IMPORTANT: Generates at most 2 slides in total that require images (meaning the sum of "visual" + "hero" slides must be <= 2) to maintain a fast and high-quality generation.

Respond ONLY with this JSON:
{{
    "presentation_title": "Concise Title",
    "slides": [
        {{
            "type": "content",
            "title": "Slide Title",
            "points": ["Point 1.", "Point 2."]
        }},
        {{
            "type": "visual",
            "title": "Slide Title",
            "points": ["Point 1.", "Point 2."],
            "image_prompt": "short concrete image description"
        }},
        {{
            "type": "comparison",
            "title": "Slide Title",
            "left_points": ["Left 1.", "Left 2."],
            "right_points": ["Right 1.", "Right 2."]
        }},
        {{
            "type": "quote",
            "quote": "The quote text.",
            "author": "Author Name"
        }},
        {{
            "type": "hero",
            "title": "Transition or Opener Title",
            "subtitle": "Supporting subtitle message.",
            "image_prompt": "short concrete image description"
        }},
        {{
            "type": "stats",
            "title": "Key Statistics or Performance Metrics",
            "stats": [
                {{"value": "95%", "label": "Satisfied Customers"}},
                {{"value": "10x", "label": "Growth Rate"}},
                {{"value": "24/7", "label": "Support Available"}}
            ]
        }}
    ]
}}
'''

    response = await generate_json_response(
        system_prompt=system_prompt,
        user_prompt=f"Generate {num_slides} highly detailed slides now.",
        model=model,
        max_tokens=6000
    )
    slides_data = response.get("slides", [])
    pres_title = response.get("presentation_title", prompt_text[:60])

    if not slides_data:
        raise ValueError("AI failed to generate structural slide data.")

    # ── Fetch visual/hero slide images sequentially (Pollinations rate-limits to 1 queued request per IP) ──
    img_paths = {}
    async with httpx.AsyncClient(verify=False) as client:
        for idx, slide in enumerate(slides_data):
            if slide.get("type") in ("visual", "hero") and slide.get("image_prompt"):
                print(f"[PPT Image] Fetching image {idx+1}: {slide['image_prompt'][:50]}...")
                path = await _fetch_image_for_slide_async(client, slide["image_prompt"])
                if path:
                    img_paths[idx] = path
                # Small delay between requests to avoid rate limiting
                await asyncio.sleep(1)

    try:
        # ── Build the presentation ──
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # 1. Title slide
        subtitle = f"{len(slides_data)} slides  •  AI-generated study material"
        _build_title_slide(prs, pres_title, subtitle)

        # 2. Content slides
        for idx, slide_data in enumerate(slides_data, start=1):
            slide_type = slide_data.get("type", "content")
            if slide_type == "visual":
                pre_fetched_img = img_paths.get(idx - 1)
                _build_visual_slide(prs, slide_data, idx, len(slides_data), image_path=pre_fetched_img)
            elif slide_type == "hero":
                pre_fetched_img = img_paths.get(idx - 1)
                _build_hero_slide(prs, slide_data, idx, len(slides_data), image_path=pre_fetched_img)
            elif slide_type == "comparison":
                _build_comparison_slide(prs, slide_data, idx, len(slides_data))
            elif slide_type == "quote":
                _build_quote_slide(prs, slide_data, idx, len(slides_data))
            elif slide_type == "stats":
                _build_stats_slide(prs, slide_data, idx, len(slides_data))
            else:
                _build_content_slide(prs, slide_data, idx, len(slides_data))

        # 3. Closing slide
        _build_closing_slide(prs, pres_title)
    finally:
        # Clean up any image files that were downloaded but not removed
        for path in img_paths.values():
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass

    # ── Save ──
    # Absolute path to output directory in backend folder
    BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    OUTPUT_DIR = os.path.join(BASE_DIR, "output")
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    # Sanitize filename
    safe_name = re.sub(r'[^a-zA-Z0-9_\- ]', '', pres_title)[:40].strip().replace(' ', '_')
    filename = f"presentation_{safe_name}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)

    # Supabase Tracking (skip for guest/anonymous users)
    if user_id not in ("guest", "anonymous"):
        from app.utils.database import log_generation
        log_generation(
            user_id=user_id, 
            content_type="ppt_presentation", 
            content_data={"prompt": prompt_text, "num_slides": len(slides_data), "filename": filename}
        )
    
    return filepath
